# Predefined categories 
CATEGORIES = {
    "Financial Reports": ["annual report", "quarterly report", "earnings"],
    "Investor Presentations": ["presentation", "conference", "slides"],
    "Corporate Governance Documents": ["policy", "charter", "governance"],
    "Press Releases": ["announcement", "merger", "leadership"],
    "Stock Market Information": ["stock price", "dividend", "shareholder"],
    "Corporate Social Responsibility (CSR) Reports": ["sustainability", "ESG", "community"]
}

# Import necessary libraries
import os
import requests
import pdfplumber
from unstructured.partition.auto import partition
from unstructured.partition.pdf import partition_pdf
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

import json
import uuid
import chromadb

# Helper function to download files from URLs and save them to a specified directory
def download_file(url, save_dir="downloads"):
    
    os.makedirs(save_dir, exist_ok=True)
    filename = os.path.join(save_dir, url.split("/")[-1])
    
    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()
        
        with open(filename, "wb") as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        
        print(f"Downloaded: {filename}")
        return filename
    except requests.RequestException as e:
        print(f"Failed to download {url}: {e}")
        return None

# Helper function to extract text from PDF files
def parse_pdf(file_path, url_source=None):
    text_data = []
    
    elements = partition_pdf(filename=file_path)
    full_text = "\n".join([str(el) for el in elements])
    category = classify_document(full_text)
    for element in elements:
        text_data.append({"page_num": element.metadata.page_number, "content": element.text, "category": category, "url": url_source, "doc_type": "pdf"})
    
    return text_data

# Helper function to extract text from PPTX files using python-pptx
def parse_ppt(file_path, url_source=None):
    text_data = []
    prs = Presentation(file_path)

    # Extract text from all slides
    full_text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    category = classify_document(full_text)

    # Extract text from each slide and store it in a structured format with page number, content, category, and source URL
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        text_data.append({"page_num": slide_num, "content": slide_text,  "category": category, "url": url_source, "doc_type": "ppt"})
    return text_data

# Helper function to chunk text into smaller segments
def chunk_text(text, chunk_size=500, overlap=100):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=overlap)
    return splitter.split_text(text)

# Helper function to save data to a JSONL file
def save_to_json(data, output_file="output/output.json"):
    os.makedirs("output", exist_ok=True)
    with open(output_file, "w") as f:
        for entry in data:
            f.write(json.dumps(entry) + "\n")

def fetch_documents(urls):
    # Implement HTTP client with retry logic
    files = []

    for url in urls:
        file_path = download_file(url)
        if file_path:
            files.append({file_path, url})

    return files

def classify_document(text):
    # Rule-based classification using keyword matching
    for category, keywords in CATEGORIES.items():
        if any(keyword.lower() in text.lower() for keyword in keywords):
            return category
    return "Unknown"

def process_file(file_path, url_source=None):
    
    # Determine the file type based on extension
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    doc_type = "pdf" if ext == ".pdf" else "ppt" if ext in [".ppt", ".pptx"] else None
    if doc_type is None:
        print(f"Unsupported file type: {ext}")
        return "Unsupported file type."
    
    # Implement file type detection and pre-processing
    extracted_data = parse_pdf(file_path, url_source) if doc_type == "pdf" else parse_ppt(file_path, url_source)

    return extracted_data
    

def chunk_data(elements):
    # Implement semantic chunking logic
    final_data = []
    
    for entry in elements:
        for item in entry:
            category = item["category"]
            chunks = chunk_text(item["content"])

            for chunk in chunks:
                final_data.append({
                    "chunk_id": str(uuid.uuid4()),
                    "content": chunk,
                    "metadata": {
                        "source": item["url"],
                        "page_num": item["page_num"],
                        "doc_type": item["doc_type"],
                        "content_type": "text",
                        "category": category
                    }
                })
    
    return final_data


def store_in_chromadb(client, data):
    collection = client.get_or_create_collection("documents")
    for entry in data:
        collection.add(
            documents=[entry["content"]],
            metadatas=[entry["metadata"]],
            ids=[entry["chunk_id"]]
        )

def main():

    urls = [
    #    "https://view.officeapps.live.com/op/view.aspx?src=https://cdn-dynmedia-1.microsoft.com/is/content/microsoftcorp/SlidesFY25Q2",
       "https://digitalassets.tesla.com/tesla-contents/image/upload/IR/TSLA-Q4-2024-Update.pdf",
    #    "https://s2.q4cdn.com/470004039/files/doc_earnings/2025/q1/filing/10Q-Q1-2025-as-filed.pdf",
    #    "https://www.apple.com/newsroom/pdfs/fy2025-q1/FY25_Q1_Consolidated_Financial_Statements.pdf",
    #    "https://s2.q4cdn.com/470004039/files/doc_financials/2021/q4/_10-K-2021-(As-Filed).pdf"
    ]

    raw_files = fetch_documents(urls)
    processed_data = [process_file(f[0],f[1]) for f in raw_files]
    chunks = chunk_data(processed_data)
    
    # save final chunks to a database or file system (output)

    client = chromadb.PersistentClient(path="chromadb")
    store_in_chromadb(client, chunks)
    print("Data processing and storage complete.")

if __name__ == "__main__":
    main()
